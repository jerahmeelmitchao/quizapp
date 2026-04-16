from io import BytesIO

from docx import Document
from fastapi import UploadFile
from pptx import Presentation
from PyPDF2 import PdfReader


async def extract_text_from_upload(file: UploadFile, ext: str) -> str:
    content = await file.read()
    buffer = BytesIO(content)

    if ext == "docx":
        return extract_docx_text(buffer)
    if ext == "pptx":
        return extract_pptx_text(buffer)
    if ext == "pdf":
        return extract_pdf_text(buffer)

    raise ValueError("Unsupported file extension")


def extract_docx_text(buffer: BytesIO) -> str:
    doc = Document(buffer)
    parts: list[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    return "\n".join(parts)


def extract_pptx_text(buffer: BytesIO) -> str:
    prs = Presentation(buffer)
    parts: list[str] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        parts.append(f"Slide {slide_index}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    parts.append(text)

    return "\n".join(parts)


def extract_pdf_text(buffer: BytesIO) -> str:
    reader = PdfReader(buffer)
    parts: list[str] = []

    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text.strip())

    return "\n".join(parts)